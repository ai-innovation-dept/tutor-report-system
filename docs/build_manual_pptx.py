# -*- coding: utf-8 -*-
"""指導実績報告システム 操作手順書を編集可能なPowerPoint(.pptx)で生成する。
python-pptx が入った環境（backendコンテナで pip install python-pptx 後）で実行する。
  例) docker compose exec backend python /tmp/build_manual_pptx.py /tmp/操作手順書.pptx
"""
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x15, 0x80, 0x3D)
RED = RGBColor(0xB9, 0x1C, 0x1C)
AMBER = RGBColor(0xB4, 0x53, 0x09)
INK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADBG = RGBColor(0xE8, 0xEE, 0xF7)
FONT = "Meiryo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def style_run(run, size=14, bold=False, color=INK, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def slide():
    return prs.slides.add_slide(BLANK)


def title_bar(s, no, text):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.35); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if no != "":
        r = p.add_run(); r.text = f"{no}　"; style_run(r, 22, True, RGBColor(0x93, 0xC5, 0xFD))
    r = p.add_run(); r.text = text; style_run(r, 20, True, WHITE)


def box(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    return tf


def para(tf, text, size=14, bold=False, color=INK, first=False, space=5, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space); p.alignment = align
    r = p.add_run(); r.text = text
    style_run(r, size, bold, color)
    return p


def table(s, data, left, top, width, col_w, fs=11.5, header=True, row_h=0.0):
    rows, cols = len(data), len(data[0])
    g = s.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.4 * rows)).table
    g.first_row = header
    for ci, w in enumerate(col_w):
        g.columns[ci].width = Inches(w)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = g.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            is_head = header and ri == 0
            style_run(r, fs, is_head, NAVY if is_head else INK)
            if is_head:
                cell.fill.solid(); cell.fill.fore_color.rgb = HEADBG
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return g


def flow(s, top):
    steps = [("講師", "報告書を作成", BLUE), ("講師", "保護者へ承認依頼", BLUE),
             ("保護者", "承認", GREEN), ("受付", "受付承認(修正可)", NAVY),
             ("再鑑", "再鑑承認", NAVY), ("管理者", "最終承認→完了", NAVY)]
    bw, bh, gap = 1.7, 1.1, 0.34
    left = 0.55
    for i, (who, what, col) in enumerate(steps):
        x = left + i * (bw + gap)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(bw), Inches(bh))
        sh.fill.solid(); sh.fill.fore_color.rgb = WHITE; sh.line.color.rgb = col; sh.line.width = Pt(1.5)
        tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, who, 11, True, col, first=True, space=2, align=PP_ALIGN.CENTER)
        para(tf, what, 11, False, INK, space=0, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw - 0.02), Inches(top + bh / 2 - 0.12), Inches(gap + 0.04), Inches(0.24))
            ar.fill.solid(); ar.fill.fore_color.rgb = BLUE; ar.line.fill.background()
    note = box(s, left, top + bh + 0.15, 11.0, 0.5)
    para(note, "※ どの段階でも「差戻し」で前の人に戻せます（差戻しには理由コメントが必要）。", 12, True, RED, first=True)


# ============ 1. タイトル ============
s = slide()
t = box(s, 1.0, 2.4, 11.3, 3)
para(t, "指導実績報告システム", 40, True, NAVY, first=True, align=PP_ALIGN.CENTER, space=10)
para(t, "操作手順書（プロトタイプお試し版）", 20, False, GREY, align=PP_ALIGN.CENTER, space=6)
para(t, "はじめての方でも、この手順どおりに進めれば一人で操作できます。", 14, False, GREY, align=PP_ALIGN.CENTER)

# ============ 2. はじめに（役割） ============
s = slide(); title_bar(s, "0", "はじめに（このシステムでできること）")
t = box(s, 0.5, 1.0, 12.3, 0.8)
para(t, "先生（講師）が毎月の指導内容を記録し、保護者と運営スタッフが順番に確認・承認していくシステムです。", 14, first=True)
table(s, [
    ["役割", "おもな役目"],
    ["管理者／管理責任者", "利用者（アカウント）の作成や、講師・保護者・生徒の紐づけを管理します。"],
    ["講師", "指導した内容の報告書を作成し、保護者へ承認をお願いします。"],
    ["保護者", "わが子の報告書を確認し、承認（または差戻し）します。"],
    ["受付 → 再鑑 → 管理者", "運営スタッフが3段階で内容を確認・承認します。修正できるのは受付だけです。"],
], 0.5, 1.9, 12.3, [2.8, 9.5], fs=13)
t2 = box(s, 0.5, 4.7, 12.3, 1)
para(t2, "ヒント：自分がどの役割かは、ログイン後に画面左のメニュー（ダッシュボード／報告書一覧／承認管理 など）で分かります。", 12, True, BLUE, first=True)

# ============ 3. 全体の流れ＋ログイン ============
s = slide(); title_bar(s, "", "全体の流れ と ログイン方法")
flow(s, 1.15)
t = box(s, 0.5, 3.7, 12.3, 3)
para(t, "● ログインのしかた（全員共通）", 15, True, NAVY, first=True, space=6)
para(t, "1. ブラウザ（Chromeなど）で、ご案内されたURLを開きます（例：http://52.197.43.164:8000）。", 13)
para(t, "2. 「メールアドレス」と「パスワード」を入力します。", 13)
para(t, "3. 「ログイン」ボタンを押します。", 13)
para(t, "ヒント：お試し用アカウントは最後のページの一覧をご覧ください（パスワードは共通で Passw0rd!）。", 12, True, BLUE, space=0)

# ============ 4. ユーザ作成 ============
s = slide(); title_bar(s, "1", "ユーザ（アカウント）を作成する　※管理者の操作")
t = box(s, 0.5, 1.0, 12.3, 5.8)
para(t, "新しく講師・保護者・運営スタッフを使えるようにするには、まず管理者がアカウントを「招待」します。", 13, first=True, space=8)
for ln in [
    "1. 管理者でログインし、左メニューの「ユーザー管理」を開きます。",
    "2. 「新規ユーザー登録」で、作りたい役割（講師／保護者／受付／再鑑／管理者）を選びます。",
    "3. 講師・運営スタッフ → 「氏名」と「メールアドレス」を入力します。",
    "4. 保護者 → 「担当講師」を選び、「生徒名」と「メールアドレス」を入力します。",
    "5. 「招待メールを送る」を押すと、相手に登録用のメールが届きます。",
    "6. メールを受け取った本人がリンクを開き、パスワードを設定すると登録完了です。",
]:
    para(t, ln, 13)
para(t, "注意：招待リンクの有効期限は72時間です。期限切れは一覧の「再送」でもう一度送れます。", 12, True, AMBER, space=4)
para(t, "ヒント：お試し環境では、送信メールを「メール確認画面（MailHog）」で確認できる場合があります（例：…:8025）。", 12, True, BLUE, space=0)

# ============ 5. ユーザー詳細の項目説明 ============
s = slide(); title_bar(s, "1", "ユーザー詳細画面の見かた（項目の説明）")
table(s, [
    ["項目", "説明"],
    ["基本情報", "No（利用者番号）・名前・メール・登録日・最終ログイン日を表示します。"],
    ["ロール設定", "役割を表示。受付／再鑑は両方を兼ねる設定をチェックで切り替え「保存」で確定します。"],
    ["担当生徒（講師の場合）", "その講師が担当する生徒の一覧を表示します。"],
    ["紐づく生徒（保護者の場合）", "その保護者に紐づく生徒の一覧。下の欄から「保護者未設定の生徒」を選んで紐づけられます。"],
    ["保護者承認をスキップ（管理責任者のみ）", "ONにすると、その保護者の承認を省略して直接運営へ提出します。"],
    ["状態管理", "「無効化する／有効化する」で、その人の利用可否を切り替えます。"],
    ["危険な操作（削除）", "ユーザーを完全に削除します。確認のためメールアドレスの入力が必要です（元に戻せません）。"],
], 0.5, 1.05, 12.3, [3.6, 8.7], fs=12.5)

# ============ 6. 紐づき設定 ============
s = slide(); title_bar(s, "2", "紐づき設定（講師・保護者・生徒）　※管理者の操作")
t = box(s, 0.5, 1.0, 12.3, 2.4)
para(t, "「誰の保護者の、どの生徒を、どの講師が担当するか」を結びつける設定です。これが無いと講師は承認依頼ができません。", 13, first=True, space=8)
para(t, "● 基本：保護者を「担当講師＋生徒名」つきで招待 → 保護者が登録すると自動でつながります（既存保護者なら即時）。", 13, True, NAVY)
para(t, "● 後から紐づけ：ユーザー管理 → 保護者の「詳細」→「紐づく生徒」→「保護者未設定の生徒」を選んで「紐づける」。", 13)
para(t, "● 担当の変更・整理は左メニューの「担当管理」で行います（下表）。", 13)
table(s, [
    ["担当管理でできること", "説明"],
    ["講師を交代／追加", "担当の先生を変更、または同じ生徒に別の講師を追加します。"],
    ["保護者を変更／解除", "保護者を選び直します。「未設定」で解除します。"],
    ["生徒名を編集", "名前の打ち間違いなどを修正します。"],
    ["有効／無効・削除", "使わない担当を無効化、または完全に削除します（報告書がある担当は削除不可）。"],
], 0.5, 3.6, 12.3, [3.2, 9.1], fs=12.5)

# ============ 7. 報告書作成 ============
s = slide(); title_bar(s, "3", "講師が報告書を作成する　※講師の操作")
t = box(s, 0.5, 1.0, 7.3, 5.8)
for ln in [
    "1. 講師でログインし、左メニューの「報告書一覧」を開きます。",
    "2. 「簡易作成」で、まず対象の「生徒」を選びます。",
    "3. 「指導日」を選びます（曜日は自動表示）。",
    "4. 「開始時刻」「終了時刻」を選びます（5分単位）。",
    "5. 「休憩等の時間（分）」を入れます（5分単位）。",
    "6. 「指導時間数」は自動計算。0.5時間（30分）単位になるよう休憩で調整。",
    "7. 「科目」「指導内容」を入力し「保存」を押します。",
]:
    para(t, ln, 13, first=(ln.startswith("1.")))
para(t, "注意：指導時間数が0.5時間（30分）単位でないと保存できません。", 12, True, AMBER, space=0)
img = box(s, 8.1, 1.05, 4.7, 4.5)
para(img, "［ 簡易作成（イメージ）］", 12, True, GREY, first=True, space=6)
for ln in ["生徒：[ 生徒を選択 ▼ ]", "指導日：[ 2026/06/10 ]（水）", "開始：[ 16:00 ] 〜 終了：[ 18:00 ]",
           "休憩（分）：[ 0 ]", "指導時間数：2時間（自動）", "科目：[ 数学 ]", "指導内容：[ ……… ]", "［ 保存 ］"]:
    para(img, ln, 12, color=GREY)

# ============ 8. 承認依頼 ============
s = slide(); title_bar(s, "4", "保護者へ承認依頼する　※講師の操作")
t = box(s, 0.5, 1.0, 12.3, 5.5)
for ln in ["1. 講師でログインし、左メニューの「承認管理」を開きます。",
           "2. 当月の報告書がまとまって表示されます。内容を確認します。",
           "3. 「保護者へまとめて承認依頼」を押すと、保護者に確認のお願いが届きます。"]:
    para(t, ln, 14, first=(ln.startswith("1.")))
para(t, "注意：「保護者が未設定のため承認依頼できません」と出たら、まだ生徒に保護者がついていません。"
        "手順2（紐づき設定）で保護者をつないでからお試しください。", 13, True, AMBER, space=6)
para(t, "ヒント：差戻された場合は、報告書を直したあと「修正済みの報告書を保護者へまとめて再依頼」で再度お願いできます。", 13, True, BLUE)

# ============ 9. 確認・承認・差戻し ============
s = slide(); title_bar(s, "5", "確認・承認・差戻し（保護者／運営）")
t = box(s, 0.5, 1.0, 12.3, 5.8)
para(t, "● 保護者の操作", 15, True, NAVY, first=True, space=4)
for ln in ["1. 保護者でログインし「承認管理」を開きます。",
           "2. 「報告書を確認」で内容を見ます。",
           "3. 問題なければ「すべて承認する」→ 確認画面で「承認は保護者本人が操作しています」にチェックして承認。",
           "4. 直してほしいときは「差戻す」を押し、理由（コメント）を入力。"]:
    para(t, ln, 13)
para(t, "● 運営スタッフの操作（受付 → 再鑑 → 管理者）", 15, True, NAVY, space=4)
for ln in ["1. 「ダッシュボード」で自分の担当の報告書を「報告書を確認」で確認。",
           "2. 問題なければ承認（受付承認／再鑑承認／最終承認）。",
           "3. 直してほしいときは「差戻し」＋理由を入力。",
           "4. 報告書の内容を修正できるのは「受付」だけ →「報告書の修正（受付）」→「修正して通知」。"]:
    para(t, ln, 13)
para(t, "ポイント：承認の順番は 講師 → 保護者 → 受付 → 再鑑 → 管理者（最終）。最後まで承認されると完了です。", 12, True, GREEN, space=0)

# ============ 10. ダッシュボードの見かた ============
s = slide(); title_bar(s, "5", "ダッシュボードの見かた（運営）")
table(s, [
    ["画面の要素", "説明"],
    ["上部のしぼり込み", "「対象月」や「講師」で表示をしぼり込めます。"],
    ["状況の列", "報告書が今どの段階か（受付待ち／再鑑待ち／最終承認待ち／完了 など）が列で並びます。"],
    ["自分の担当タスク", "自分が対応する段階の報告書に、承認ボタンと「差戻し」が表示されます。"],
    ["報告書を確認", "指導日・時間・科目などの内容を別画面で確認できます。"],
    ["承認／差戻し", "承認ボタン（受付承認／再鑑承認／最終承認）または「差戻し」を押します。"],
], 0.5, 1.05, 12.3, [3.0, 9.3], fs=12.5)
t = box(s, 0.5, 4.6, 12.3, 1)
para(t, "注意：公正のため、自分が「受付」した報告書は同じ人が「再鑑」できません（ボタンが押せません）。", 13, True, AMBER, first=True)

# ============ 11. 困ったとき ============
s = slide(); title_bar(s, "6", "困ったとき（よくある質問）")
table(s, [
    ["こんなとき", "対処"],
    ["ログインできない", "メール・パスワードを確認。忘れた場合はログイン画面「パスワードをお忘れの方はこちら」から再設定。"],
    ["承認依頼ボタンが押せない／出ない", "生徒に保護者がついていない可能性。手順2で紐づけてください（当月分のみ操作可）。"],
    ["保存できない（報告書）", "指導時間数が0.5時間（30分）単位になっていません。休憩で調整してください。"],
    ["差戻しできない", "差戻しには理由（コメント）の入力が必要です。"],
    ["同じ人が受付と再鑑を両方できない", "公正のため、同じ報告書では受付した人は再鑑できません（逆も同様）。"],
], 0.5, 1.05, 12.3, [3.6, 8.7], fs=12.5)

# ============ 12. ルール ============
s = slide(); title_bar(s, "7", "当システムのルール（大切な約束）")
t = box(s, 0.5, 1.0, 12.3, 5.8)
for i, r in enumerate([
    "承認は決まった順番：講師 → 保護者 → 受付 → 再鑑 → 管理者（最終承認）。",
    "差戻し（修正のお願い）には必ず理由（コメント）を入力。差し戻すと一つ前の人へ戻ります。",
    "生徒に保護者がついていないと、講師は承認依頼ができません（先に紐づけが必要）。",
    "報告書の時刻は5分単位、指導時間数は0.5時間（30分）単位で記録します。",
    "報告書の内容を修正できるのは「受付」だけ（修正すると講師・保護者へ通知）。",
    "公正のため、同じ報告書で「受付」と「再鑑」を同じ人が兼ねることはできません。",
    "操作できるのは当月分の報告書です（過去月は参照のみ）。",
    "招待メールのリンクは72時間で期限切れになります。",
], 1):
    para(t, f"{i}.　{r}", 14, first=(i == 1), space=7)

# ============ 13. デモアカウント ============
s = slide(); title_bar(s, "", "付録：デモ用アカウント一覧（お試し用）")
table(s, [
    ["役割", "メールアドレス", "できること"],
    ["管理者", "master1@example.com", "ユーザ作成・紐づけ・最終承認"],
    ["講師", "tutor1@example.com", "報告書の作成・承認依頼"],
    ["保護者", "parent1@example.com", "報告書の承認・差戻し"],
    ["受付", "receiver1@example.com", "受付承認・修正・差戻し"],
    ["再鑑", "reviewer1@example.com", "再鑑承認・差戻し"],
], 0.5, 1.05, 12.3, [2.2, 5.0, 5.1], fs=13)
t = box(s, 0.5, 4.4, 12.3, 1.6)
para(t, "※ パスワードはすべて共通：Passw0rd!（お試し環境用）", 13, True, INK, first=True, space=8)
para(t, "おすすめのお試し順：①管理者で紐づけ確認 → ②講師で報告書作成・承認依頼 → ③保護者で承認 → "
        "④受付→再鑑→管理者の順に承認（「完了」まで体験）。", 13, True, BLUE)

out = sys.argv[1] if len(sys.argv) > 1 else "操作手順書.pptx"
prs.save(out)
print("wrote", out)
